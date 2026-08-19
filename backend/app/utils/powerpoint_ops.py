"""
PDF -> PowerPoint conversion backing /api/pdf/to-powerpoint.

There's no well-maintained library that reconstructs arbitrary PDF content
into genuinely editable PowerPoint text boxes the way `pdf2docx` does for
Word — faithfully preserving free-form layout as editable shapes is a much
harder problem for slides than for word-processing documents (which are
just flowing paragraphs). Instead, each PDF page is rendered to a
full-resolution image and placed full-bleed on its own slide — the same
approach other PDF-to-PPT tools (iLovePDF included) use under the hood.
This preserves the exact visual layout, works uniformly for any PDF
(scanned or not — no OCR needed, unlike PDF to Word/Markdown), and the
result can still be presented/annotated in PowerPoint; it just isn't
editable text.
"""
import io

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Emu

EMU_PER_POINT = 12700  # 914400 EMU/inch / 72 points/inch
RENDER_DPI = 150

# python-pptx's built-in "Blank" slide layout (index 6 in the default
# template) — no placeholder text boxes to work around.
_BLANK_LAYOUT_INDEX = 6


class PowerPointToolError(RuntimeError):
    """Raised when a PDF can't be converted to a PowerPoint — callers turn
    this into a 400."""


def pdf_to_pptx(file_bytes: bytes) -> bytes:
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as exc:
        raise PowerPointToolError(f"Couldn't read this file as a PDF: {exc}") from exc

    try:
        if doc.page_count == 0:
            raise PowerPointToolError("This PDF has no pages.")

        presentation = Presentation()
        # PowerPoint has one canvas size for the whole deck — sized to the
        # first page; any differently-sized page is scaled to fit within
        # it (preserving aspect ratio, centered) rather than stretched.
        canvas_width_emu = round(doc[0].rect.width * EMU_PER_POINT)
        canvas_height_emu = round(doc[0].rect.height * EMU_PER_POINT)
        presentation.slide_width = canvas_width_emu
        presentation.slide_height = canvas_height_emu
        blank_layout = presentation.slide_layouts[_BLANK_LAYOUT_INDEX]

        for page in doc:
            page_width_emu = round(page.rect.width * EMU_PER_POINT)
            page_height_emu = round(page.rect.height * EMU_PER_POINT)

            image_bytes = page.get_pixmap(dpi=RENDER_DPI).tobytes("png")
            slide = presentation.slides.add_slide(blank_layout)

            if page_width_emu == canvas_width_emu and page_height_emu == canvas_height_emu:
                left, top = 0, 0
                draw_width, draw_height = canvas_width_emu, canvas_height_emu
            else:
                scale = min(
                    canvas_width_emu / page_width_emu, canvas_height_emu / page_height_emu
                )
                draw_width = round(page_width_emu * scale)
                draw_height = round(page_height_emu * scale)
                left = (canvas_width_emu - draw_width) // 2
                top = (canvas_height_emu - draw_height) // 2

            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                Emu(left),
                Emu(top),
                width=Emu(draw_width),
                height=Emu(draw_height),
            )
    except PowerPointToolError:
        raise
    except Exception as exc:
        raise PowerPointToolError(f"PDF to PowerPoint conversion failed: {exc}") from exc
    finally:
        doc.close()

    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()
